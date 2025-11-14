 #Set-ExecutionPolicy -Scope Process -ExecutionPolicy Bypass -Force; .\.venv\Scripts\Activate

# Clean RAG Starter (LangChain 0.2 + Chroma + Gradio)

import os
import json
import hashlib
from typing import Dict, List

from dotenv import load_dotenv
#import chromadb
import tiktoken
load_dotenv()
#-----DO-----------

import psycopg
from langchain_postgres import PGVector
from urllib.parse import quote_plus

PG_USER = os.getenv("PG_USER", "doadmin")
PG_PASSWORD = os.getenv("PG_PASSWORD")          # put raw password in .env, NOT encoded
PG_HOST = os.getenv("PG_HOST")                  # emmett-...ondigitalocean.com (or private-... if same VPC)
PG_PORT = os.getenv("PG_PORT", "25060")
PG_DB = os.getenv("PG_DATABASE", "emmett_ai")

ENC_PWD = quote_plus(PG_PASSWORD)
PG_DSN = f"postgresql://{PG_USER}:{ENC_PWD}@{PG_HOST}:{PG_PORT}/{PG_DB}?sslmode=require"
PG_SA_URL = f"postgresql+psycopg://{PG_USER}:{ENC_PWD}@{PG_HOST}:{PG_PORT}/{PG_DB}?sslmode=require"

os.environ["PG_CONN_STR"] = PG_SA_URL

# ---- LangChain v0.2+ imports ----
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.output_parsers import StrOutputParser
from langchain_core.chat_history import InMemoryChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_core.runnables import RunnablePassthrough
from operator import itemgetter

from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_chroma import Chroma

# Token/cost tracking

from langchain_community.callbacks import get_openai_callback


from langchain_community.retrievers import BM25Retriever
#from langchain.retrievers import EnsembleRetriever

import logging
logging.getLogger("pypdf").setLevel(logging.ERROR)
logging.getLogger("PyPDF2").setLevel(logging.ERROR)

PG_CONN_STR="postgresql://doadmin:AVNS_y-18CdRh05UXrX-Z7ib@emmett-green-gis-do-user-16508034-0.c.db.ondigitalocean.com:25060/emmett_ai?sslmode=require"
# ---------------- Config ----------------
with psycopg.connect(PG_DSN) as conn:
    with conn.cursor() as cur:
        cur.execute("SELECT version()")
        print(cur.fetchone())

MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
EMBEDDING_MODEL = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-large")
EMBED_PRICE_PER_1M = float(os.getenv("EMBED_PRICE_PER_1M", "0.13"))
DB_DIR = os.getenv("DB_DIR", "vector_db")
COLLECTION = os.getenv("COLLECTION", "docs")
BASE_DIR = os.getenv("BASE_DIR", r"data")
MANIFEST_NAME = "manifest.json"
TOP_K = int(os.getenv("TOP_K", "5"))

os.makedirs(BASE_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

embeddings = OpenAIEmbeddings(model=EMBEDDING_MODEL)
from langchain_postgres import PGVector
from langchain_openai import OpenAIEmbeddings
import os

MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
EMBEDDING_MODEL = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-large")

embeddings = OpenAIEmbeddings(model=EMBEDDING_MODEL)

vectorstore = PGVector(
    embeddings=embeddings,
    collection_name=os.getenv("COLLECTION", "docs"),
    connection=PG_SA_URL,
    use_jsonb=True,
    create_extension=True,  # keep False if you created pgvector yourself
)

BM25_CORPUS: List[Document] = []


def make_pg_vectorstore():
    return PGVector(
        embeddings=embeddings,
        collection_name=os.getenv("COLLECTION", "docs"),
        connection=PG_SA_URL,
        use_jsonb=True,
        create_extension=False,
    )
# -------------- Token counting--------------

MODEL_PRICES = { "gpt-4o-mini": {"input": 0.15, "output": 0.60, "cached": 0.075}, }

def compute_cost(prompt_tokens, completion_tokens, model=MODEL, cached_fraction=0.0):
    """Compute exact cost with cached input discount."""
    p = MODEL_PRICES.get(model, MODEL_PRICES["gpt-4o-mini"])
    cached_tokens = int(prompt_tokens * cached_fraction)
    normal_tokens = prompt_tokens - cached_tokens
    cost_input = (normal_tokens * p["input"] + cached_tokens * p["cached"]) / 1000000
    cost_output = (completion_tokens * p["output"]) / 1000000
    return cost_input + cost_output

def _count_tokens_texts(texts, model_name: str) -> int:
    try:
        enc = tiktoken.encoding_for_model(model_name)
    except Exception:
        enc = tiktoken.get_encoding("cl100k_base")
    return sum(len(enc.encode(t or "")) for t in texts)

# ---------- Helpers for tabular and pptx ----------

def _normalize_ws(s: str) -> str:
    return " ".join(str(s).split()) if s is not None else ""

def _df_block_to_text(df):
    """Render a DataFrame block (no trimming)."""
    import pandas as pd
    if not isinstance(df, pd.DataFrame) or df.empty:
        return ""
    cols_line = " | ".join(map(str, df.columns.tolist()))
    lines = [cols_line]
    append = lines.append
    for _, row in df.iterrows():
        append(" | ".join("" if pd.isna(v) else _normalize_ws(v) for v in row.tolist()))
    return "\n".join(lines)

class ExcelStreamingLoader:
    """
    Streams ALL cells from ALL sheets using openpyxl read_only mode.
    - No header assumptions.
    - No row/column restrictions (everything emitted).
    - Emits Documents per N rows to avoid huge single strings.
    """
    def __init__(self, path: str, rows_per_chunk: int = 1000):
        self.path = path
        self.rows_per_chunk = rows_per_chunk

    @staticmethod
    def _safe_str(v):
        # Normalize any value to a clean string (handles None, dates, numbers, etc.)
        import datetime as _dt
        if v is None:
            return ""
        if isinstance(v, (_dt.datetime, _dt.date, _dt.time)):
            try:
                return v.isoformat()
            except Exception:
                return str(v)
        return str(v)

    def _rows_to_text(self, rows):
        # Serialize rows as " | "-separated lines (no header, no trimming)
        out = []
        append = out.append
        for r in rows:
            append(" | ".join(self._safe_str(c) for c in r))
        return "\n".join(out)

    def load(self):
        from openpyxl import load_workbook
        from langchain_core.documents import Document

        # read_only=True => streaming; data_only=True => values (not formulas)
        wb = load_workbook(self.path, read_only=True, data_only=True)
        docs = []

        for ws in wb.worksheets:
            rows_iter = ws.iter_rows(values_only=True)
            buffer, count, chunk_id, total_rows = [], 0, 0, 0

            for row in rows_iter:
                buffer.append(row)
                count += 1
                total_rows += 1

                if count >= self.rows_per_chunk:
                    text = (
                        f"[Excel sheet: {ws.title} | rows {total_rows - count + 1}-{total_rows}]\n"
                        + self._rows_to_text(buffer)
                    )
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={
                                "source": self.path,
                                "sheet": ws.title,
                                "chunk_id": chunk_id,
                                "rows_in_chunk": count,
                                "rows_total_seen": total_rows,
                                "no_split": True,
                            },
                        )
                    )
                    buffer, count, chunk_id = [], 0, chunk_id + 1

            # tail
            if buffer:
                text = (
                    f"[Excel sheet: {ws.title} | rows {total_rows - count + 1}-{total_rows}]\n"
                    + self._rows_to_text(buffer)
                )
                docs.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": self.path,
                            "sheet": ws.title,
                            "chunk_id": chunk_id,
                            "rows_in_chunk": count,
                            "rows_total_seen": total_rows,
                            "no_split": True,
                        },
                    )
                )

            # Explicitly mark truly empty sheets
            if total_rows == 0:
                docs.append(
                    Document(
                        page_content=f"[Excel sheet: {ws.title}] (empty)",
                        metadata={"source": self.path, "sheet": ws.title, "rows_total": 0, "no_split": True},
                    )
                )

        return docs


class CSVChunkedLoader:
    """
    Reads ALL rows, ALL columns from CSV with streaming.
    Uses pandas read_csv(..., chunksize=...) to avoid memory blowups.
    Emits one Document per chunk—no trimming.
    """
    def __init__(self, path: str, chunksize: int = 20000, encoding: str | None = None):
        self.path = path
        self.chunksize = chunksize  # not a limit—just splits into many docs
        self.encoding = encoding

    def load(self):
        import pandas as pd
        from langchain_core.documents import Document

        docs = []
        chunk_iter = pd.read_csv(self.path, dtype=object, chunksize=self.chunksize, encoding=self.encoding, low_memory=False)
        total_rows = 0
        chunk_id = 0
        for chunk in chunk_iter:
            rows = len(chunk)
            total_rows += rows
            text = f"[CSV chunk {chunk_id} | rows {total_rows-rows}-{total_rows-1}]\n{_df_block_to_text(chunk)}"
            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": self.path,
                        "chunk_id": chunk_id,
                        "rows_in_chunk": rows,
                        "rows_total_seen": total_rows,
                    },
                )
            )
            chunk_id += 1

        if not docs:  # empty file
            docs.append(
                Document(
                    page_content="[CSV] (empty file)",
                    metadata={"source": self.path, "rows_total_seen": 0},
                )
            )
        else:
            # annotate total rows on the last doc
            docs[-1].metadata["rows_total_final"] = total_rows
        return docs

def _table_to_text(table) -> str:
    """Extract text from python-pptx table."""
    rows = []
    for r in table.rows:
        cells = []
        for c in r.cells:
            cells.append(_normalize_ws(c.text))
        rows.append(" | ".join(cells))
    return "\n".join(rows)

class PptxRichLoader:
    """
    Extracts ALL textual content from slides (titles, shapes, tables, notes).
    Emits one Document per slide—no trimming.
    """
    def __init__(self, path: str):
        self.path = path

    def load(self):
        from pptx import Presentation
        from langchain_core.documents import Document

        prs = Presentation(self.path)
        docs = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts = []
            # Title placeholder
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    parts.append(f"[Title] {slide.shapes.title.text}")
            except Exception:
                pass
            # All shapes (text & tables)
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    txt = shape.text or ""
                    if txt.strip():
                        parts.append(txt)
                if hasattr(shape, "has_table") and shape.has_table:
                    parts.append("[Table]\n" + _table_to_text(shape.table))
            # Notes
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    parts.append("[Notes]\n" + notes)

            content = "\n".join(p for p in (s.strip() for s in parts) if p)
            docs.append(
                Document(
                    page_content=f"[Slide {idx}/{len(prs.slides)}]\n{content}",
                    metadata={"source": self.path, "slide_number": idx, "slides_total": len(prs.slides)},
                )
            )
        if not docs:
            docs.append(
                Document(
                    page_content="[PowerPoint] (no readable text found)",
                    metadata={"source": self.path, "slides_total": 0},
                )
            )
        return docs


# -------------- Helpers --------------

def sha256_file(path: str, chunk_size: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            b = f.read(chunk_size)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def pick_loader(path: str):
    """
    Zero-truncation policy:
    - Excel: all rows/cols (split into many Documents by row blocks)
    - CSV: all rows/cols (streamed into many Documents)
    - PowerPoint: all text (titles, shapes, tables, notes) per slide
    - PDF/DOCX/TXT: unchanged
    """
    from langchain_community.document_loaders import TextLoader, PyPDFLoader, Docx2txtLoader
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return PyPDFLoader(path)

    if ext == ".docx":
        return Docx2txtLoader(path)

    if ext == ".txt":
        return TextLoader(path, autodetect_encoding=True)

    if ext in {".csv"}:
        # If you often see messy encodings, try encoding=None to let pandas detect.
        return CSVChunkedLoader(path, chunksize=20000, encoding=None)

    if ext in {".xlsx", ".xls"}:
        # Adjust rows_per_chunk for your environment; it does NOT drop data.
        return ExcelStreamingLoader(path, rows_per_chunk=800)  # keep entire file, just streamed

    if ext in {".pptx", ".ppt"}:
        return PptxRichLoader(path)

    # Fallback
    return TextLoader(path, autodetect_encoding=True)


def load_manifest(db_dir: str) -> Dict[str, Dict]:
    mpath = os.path.join(db_dir, MANIFEST_NAME)
    if os.path.exists(mpath):
        with open(mpath, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_manifest(db_dir: str, manifest: Dict[str, Dict]) -> None:
    os.makedirs(db_dir, exist_ok=True)
    mpath = os.path.join(db_dir, MANIFEST_NAME)
    with open(mpath, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)


def scan_files(base_dir: str, exts=(".pdf", ".docx", ".txt", ".csv", ".xlsx", ".xls", ".pptx",".ppt")) -> List[str]:
    paths = []
    for root, _, files in os.walk(base_dir):
        for fn in files:
            if fn.lower().endswith(exts):
                paths.append(os.path.join(root, fn))
    return sorted(paths)


def split_docs(docs: List[Document]) -> List[Document]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=300,  # tokens
        chunk_overlap=60,
        encoding_name="cl100k_base"
    )
    return splitter.split_documents(docs)
"""
OLD
 splitter = RecursiveCharacterTextSplitter(
        chunk_size=400, chunk_overlap=60, separators=["\n\n", "\n", " ", ""], add_start_index=True
    )
    return splitter.split_documents(docs)
"""

# -------------- Ingestion --------------

def incremental_ingest(base_dir: str, collection_name: str, embeddings, dry_run: bool = False):
    """
    Incremental ingest into PGVector (Managed Postgres).
    - Keeps your manifest + hashing.
    - Deletes old chunks by metadata filter: {"source": <path>}
    - Tracks embedding token cost using the OpenAI callback.
    """
    vectordb = make_pg_vectorstore()
    bm25_corpus_local: List[Document] = []

    manifest = load_manifest(".")  # manifest file will live next to app; change path if you prefer
    current_files = scan_files(base_dir)

    changed, deleted, skipped = [], [], []
    current_hashes = {}
    for p in current_files:
        try:
            current_hashes[p] = sha256_file(p)
        except Exception as e:
            print(f"[WARN] Skipping unreadable file: {p} ({e})")

    # deletions
    for p in list(manifest.keys()):
        if p not in current_hashes:
            print(f"[DEL] {p}")
            if not dry_run:
                try:
                    vectordb.delete(filter={"source": p})
                except Exception as e:
                    print(f"[WARN] delete failed for {p}: {e}")
            manifest.pop(p, None)
            deleted.append(p)

    # token/cost accumulator for embeddings
    embed_usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0, "total_cost": 0.0}

    # new/changed
    for p, h in current_hashes.items():
        prev = manifest.get(p, {}).get("hash")
        if prev == h:
            try:
                loader = pick_loader(p)
                file_docs = loader.load()
                for d in file_docs:
                    d.metadata["source"] = p
                    d.metadata["file_hash"] = h
                chunks = split_docs(file_docs)
                bm25_corpus_local.extend(chunks)
            except Exception as e:
                print(f"[WARN] Could not load for BM25 (skipped file): {p} ({e})")
            skipped.append(p)
            continue

        print(f"[INDEX] {p} (new or changed)")
        if dry_run:
            continue

        # remove old chunks for this file before re-adding
        try:
            vectordb.delete(filter={"source": p})
        except Exception as e:
            print(f"[WARN] delete failed for {p}: {e}")

        loader = pick_loader(p)
        file_docs = loader.load()
        for d in file_docs:
            d.metadata["source"] = p
            d.metadata["file_hash"] = h

        chunks = split_docs(file_docs)
        bm25_corpus_local.extend(chunks)

        texts = [c.page_content for c in chunks]
        metadatas = []
        ids = []
        for i, c in enumerate(chunks):
            md = dict(c.metadata or {})
            md["chunk_id"] = i
            md["doc_id"] = f"{p}:::{i}"
            metadatas.append(md)
            ids.append(md["doc_id"])

        # This will call OpenAIEmbeddings under the hood; the callback captures token usage
        from langchain_community.callbacks import get_openai_callback
        with get_openai_callback() as cb:
            vectordb.add_texts(texts=texts, metadatas=metadatas, ids=ids)

        emb_tokens = getattr(cb, "total_embedding_tokens", 0) or getattr(cb, "total_tokens", 0)
        if emb_tokens == 0:
            emb_tokens = _count_tokens_texts(texts, EMBEDDING_MODEL)

        embed_usage["total_tokens"] += emb_tokens
        embed_usage["prompt_tokens"] += emb_tokens
        embed_usage["total_cost"] += (emb_tokens / 1_000_000.0) * EMBED_PRICE_PER_1M

        manifest[p] = {"hash": h}
        changed.append(p)

    save_manifest(".", manifest)

    # count rows (optional)
    try:
        with psycopg.connect(PG_CONN_STR) as conn, conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM langchain_pg_embedding WHERE collection_id = "
                        "(SELECT uuid FROM langchain_pg_collection WHERE name=%s);", (collection_name,))
            total_rows = cur.fetchone()[0]
    except Exception:
        total_rows = None

    summary = {
        "changed": changed,
        "deleted": deleted,
        "skipped": skipped,
        "total_in_db": total_rows,
        "embedding_usage": embed_usage,
    }
    return vectordb, bm25_corpus_local, summary


# ---- Run ingestion once on startup ----
vectorstore, BM25_CORPUS, summary = incremental_ingest(
    base_dir=BASE_DIR,
    collection_name=COLLECTION,
    embeddings=embeddings,
    dry_run=False,
)


u = summary["embedding_usage"]
print(
    f"\nIngest summary:\n"
    f"  Added/Updated files: {len(summary['changed'])}\n"
    f"  Deleted files: {len(summary['deleted'])}\n"
    f"  Unchanged files: {len(summary['skipped'])}\n"
    f"  Total chunks in DB: {summary['total_in_db']}\n"
    f"  Embedding tokens this run: {u['total_tokens']} (cost ${u['total_cost']:.6f})\n"
)

# -------------- RAG Chain --------------

# ---- Keep your rrf_fuse exactly as you wrote it ----
def rrf_fuse(result_lists, weights=None, k=TOP_K, c=60, id_key="doc_id"):
    weights = weights or [1.0] * len(result_lists)
    scores, by_id = {}, {}
    for i, docs in enumerate(result_lists):
        w = weights[i]
        for rank, d in enumerate(docs):
            did = (d.metadata or {}).get(id_key) or d.page_content
            by_id[did] = d
            scores[did] = scores.get(did, 0.0) + w * (1.0 / (c + rank + 1))
    return [by_id[did] for did, _ in sorted(scores.items(), key=lambda x: x[1], reverse=True)[:k]]



# after incremental_ingest() returns:

bm25 = BM25Retriever.from_documents(BM25_CORPUS, k=TOP_K)

# (optional) ensemble:
#from langchain_community.retrievers  import EnsembleRetriever
#dense = vectorstore.as_retriever(search_kwargs={"k": TOP_K})
dense = vectorstore.as_retriever(search_type="mmr", search_kwargs={"k": TOP_K, "fetch_k": 20, "lambda_mult": 0.7})
#retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={"k": TOP_K, "fetch_k": 20, "lambda_mult": 0.7})
#retriever = EnsembleRetriever(retrievers=[bm25, dense], weights=[0.4, 0.6])

# ---- Minimal adapter so hybrid behaves like a retriever ----
class HybridRRFRetriever:
    def __init__(self, bm25_ret, dense_ret, weights=(0.4, 0.6), k=TOP_K, c=60, id_key="doc_id"):
        self.bm25 = bm25_ret
        self.dense = dense_ret
        self.weights = list(weights)
        self.k = k
        self.c = c
        self.id_key = id_key

    def _retrieve(self, query: str):
        # Use .invoke(...) for both retrievers (works in your versions)
        bm25_docs  = self.bm25.invoke(query)
        dense_docs = self.dense.invoke(query)
        return rrf_fuse(
            [bm25_docs, dense_docs],
            weights=self.weights,
            k=self.k,
            c=self.c,
            id_key=self.id_key,
        )

    # Support both calling styles
    def invoke(self, query: str, **kwargs):
        return self._retrieve(query)

    def get_relevant_documents(self, query: str):
        return self._retrieve(query)

# Instantiate the hybrid retriever
hybrid = HybridRRFRetriever(bm25, dense, weights=(0.4, 0.6), k=TOP_K, c=60, id_key="doc_id")



llm = ChatOpenAI(temperature=0.0, top_p=0.9, model=MODEL, streaming=True, stream_usage=True)

rewrite_prompt = ChatPromptTemplate.from_messages([
    (
        "system",
        "Given the chat history and the latest user message, rewrite the message into a standalone question. "
        "Only return the rewritten question.",
    ),
    MessagesPlaceholder("chat_history"),
    ("human", "{question}"),
])


def format_docs(docs):
    formatted = []
    for d in docs:
        src = d.metadata.get("source", "unknown")
        src_name = os.path.basename(src)
        formatted.append(f"[Source: {src_name}]\n{d.page_content}")
    return "\n\n---\n\n".join(formatted)


prompt = ChatPromptTemplate.from_messages([
    ("system", "You are a helpful assistant. Use the provided context to answer accurately and concisely."),
    MessagesPlaceholder("chat_history"),
    ("system", "Context:\n{context}"),
    ("human", "{question}"),
])

contextualize_q_chain = rewrite_prompt | llm | StrOutputParser()
"""
answer_prompt = ChatPromptTemplate.from_messages([
    (
        "system",
        "You are a helpful assistant. Use ONLY the provided context to answer accurately and concisely. "
        "Retrieve data only from the database and at the end of your answer, always append a tag in the format: [File: <filename>] "
        "Always give an answer. If the conxtext is not from the database,try to answer the question but at the end append [File: None]. ",
    ),  
    ("system", "Context:\n{context}"),
    ("human", "{question}"),
])
"""
answer_prompt = ChatPromptTemplate.from_messages([
    (
        "system",
        "You write in a very detailed, elegant tone, formatting all answers in Markdown. "
        "You are a helpful assistant. If the user expresses gratitude (e.g., 'thank you' or 'thanks'), "
        "reply politely friendly response (e.g., 'You're very welcome!'). "
        "Otherwise, use ONLY the provided context to answer accurately and concisely. "
        "Retrieve data only from the database and at the end of your answer, always append a tag in the format: [File: <filenames>]. "
        "If question is not from the database try to answer but mention it is not from the database. Then append [File: None]."
    ),
    ("system", "Context:\n{context}"),
    ("human", "{question}"),
])


def maybe_rewrite(question, chat_history):
    if not chat_history or len(chat_history) < 2:
        return question  # no rewrite needed
    return contextualize_q_chain.invoke({"question": question, "chat_history": chat_history[:]})

from langchain_core.runnables import RunnableLambda, RunnablePassthrough
from operator import itemgetter

retriever_runnable = RunnableLambda(lambda q: hybrid.invoke(q))
format_runnable    = RunnableLambda(format_docs)

rag_chain = (
    RunnablePassthrough()
    .assign(standalone_question=lambda x: maybe_rewrite(x["question"], x.get("chat_history", [])))
    .assign(context=(itemgetter("standalone_question") | retriever_runnable | format_runnable))
    .assign(question=itemgetter("standalone_question"))
    | answer_prompt
    | llm
    | StrOutputParser()
)


# ---- Chat history (memory) ----
_histories = {}


def get_history(session_id: str) -> InMemoryChatMessageHistory:
    if session_id not in _histories:
        _histories[session_id] = InMemoryChatMessageHistory()
    return _histories[session_id]


chat_chain = RunnableWithMessageHistory(
    rag_chain, get_history, input_messages_key="question", history_messages_key="chat_history"
)

SESSION_ID = os.getenv("SESSION_ID", "default-session")

def is_trivial(query: str) -> bool:
    q = (query or "").strip().lower()
    return len(q.split()) <= 2 and q in {"hi","hey","hello","yo","sup","good morning","good evening"}

# ---- Add near your globals ----
from collections import defaultdict

USAGE_TOTALS = defaultdict(float)  # keys: prompt_tokens, completion_tokens, total_tokens, total_cost
SESSION_TOTALS = defaultdict(lambda: defaultdict(float))  # per-session (SESSION_ID) buckets

def chat(message, history):
    if is_trivial(message):
        return "Hi I am EMMETT.ai! How can I assist you today?"

    with get_openai_callback() as cb:
        result = chat_chain.invoke(
            {"question": message},
            config={"callbacks": [cb],  "configurable": {"session_id": SESSION_ID}},
        )

    # use actual pricing
    cost = compute_cost(cb.prompt_tokens, cb.completion_tokens, model=MODEL)

    # update accumulators
    USAGE_TOTALS["prompt_tokens"]     += cb.prompt_tokens
    USAGE_TOTALS["completion_tokens"] += cb.completion_tokens
    USAGE_TOTALS["total_tokens"]      += cb.total_tokens
    USAGE_TOTALS["total_cost"]        += cost

    # per-call and running totals together
    print(
        f"[CHAT USAGE] Prompt={cb.prompt_tokens}  Completion={cb.completion_tokens}  "
        f"Total={cb.total_tokens}  Cost=${cost:.6f}  ||  "
        f"[TOTAL USAGE] Tokens={USAGE_TOTALS['total_tokens']:.0f}  "
        f"Cost=${USAGE_TOTALS['total_cost']:.4f}"
    )

    return result



# === RAG ADAPTER: expose a callable function for the agent ===

def rag_answer(question: str, session_id: str = SESSION_ID, top_k: int = TOP_K):
    """
    Returns a structured result the agent can consume.
    - Uses your rewrite chain to get a standalone question.
    - Retrieves docs from the ensemble retriever.
    - Calls the answer LLM prompt.
    - Returns the final answer and the retrieved docs (for sources).
    """
    # 1) Make standalone question (same logic as your chain)
    standalone = maybe_rewrite(question, get_history(session_id).messages if session_id in _histories else [])

    # 2) Retrieve docs using THE STANDALONE QUESTION (fixes your earlier context-routing gap)
    docs = hybrid.invoke(standalone)

    # 3) Format context and call your answer chain components explicitly
    context_text = format_docs(docs)
    final = (answer_prompt | llm | StrOutputParser()).invoke({"context": context_text, "question": standalone})

    # 4) Build a lean, agent-friendly payload
    sources = []
    for d in docs[:top_k]:
        md = d.metadata or {}
        sources.append({
            "file": os.path.basename(md.get("source", "unknown")),
            "full_path": md.get("source", "unknown"),
            "chunk_id": md.get("chunk_id", None),
            "snippet": d.page_content[:280]
        })

    # Simple confidence heuristic: number of hits (you can refine later)
    confidence = min(1.0, len(docs) / 5.0)

    return {
        "answer": final,
        "sources": sources,
        "used_db": True,
        "confidence": confidence,
    }
from langchain_core.tools import tool
# === Tool wrapper (NEW) ===
@tool("rag_lookup")
def rag_lookup(question: str) -> str:
    """
    Retrieve an answer from the local document database.
    Returns a JSON string with: answer, sources [{file, full_path, chunk_id, snippet}], used_db, confidence.
    """
    out = rag_answer(question)
    # Return JSON (strings are best for tool outputs in ReAct loops)
    return json.dumps(out, ensure_ascii=False)
